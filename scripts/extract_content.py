#!/usr/bin/env python3
"""
extract_content.py — 从多种格式提取文本，输出统一的 Markdown。

支持格式:
    .md / .txt   — 编码规范化，复制到 extracted/
    .pdf         — 逐页提取文本，添加 ## Page N 标记
    .pptx        — 逐幻灯片提取文本（标题、文本框、备注），添加 ## Slide N 标记

用法:
    python scripts/extract_content.py --input input/chapter4.pdf
    python scripts/extract_content.py --input input/chapter4.pptx
    python scripts/extract_content.py --input input/chapter4.md
    python scripts/extract_content.py --input input/chapter4.md -o extracted/my_notes.md

输出:
    extracted/<source_stem>.md（默认）或通过 -o 指定
"""

import argparse
import sys
from abc import ABC, abstractmethod
from pathlib import Path

# 确保 scripts/ 目录可被导入
_scripts_dir = Path(__file__).resolve().parent
if str(_scripts_dir) not in sys.path:
    sys.path.insert(0, str(_scripts_dir))

PROJECT_ROOT = _scripts_dir.parent
EXTRACTED_DIR = PROJECT_ROOT / "extracted"


# ============================================================
# 提取器抽象基类
# ============================================================

class BaseExtractor(ABC):
    """提取器基类。新增格式只需实现 extract()。"""

    @abstractmethod
    def extract(self, path: Path) -> str:
        """从文件提取文本，返回 Markdown 字符串。"""
        ...

    @staticmethod
    def for_path(path: Path) -> "BaseExtractor":
        """根据文件扩展名创建对应的提取器。"""
        ext = path.suffix.lower()
        if ext == ".pdf":
            return PdfExtractor()
        elif ext == ".pptx":
            return PptxExtractor()
        elif ext in (".md", ".markdown"):
            return MarkdownExtractor()
        elif ext == ".txt":
            return TextExtractor()
        else:
            raise ValueError(
                f"不支持的文件格式: {ext}\n"
                f"  支持的格式: .pdf, .pptx, .md, .txt"
            )


# ============================================================
# Markdown / 纯文本提取器
# ============================================================

class MarkdownExtractor(BaseExtractor):
    """Markdown 文件：规范化编码后直接复制内容。"""

    def extract(self, path: Path) -> str:
        text = _read_with_fallback(path)
        if not text.strip():
            raise ValueError(f"Markdown 文件为空: {path}")
        return text


class TextExtractor(BaseExtractor):
    """纯文本文件：转为 Markdown 格式。"""

    def extract(self, path: Path) -> str:
        text = _read_with_fallback(path)
        if not text.strip():
            raise ValueError(f"文本文件为空: {path}")
        # 如果原文没有任何 # 标题，用文件名生成一级标题
        if not any(line.strip().startswith("#") for line in text.splitlines()):
            title = path.stem
            text = f"# {title}\n\n{text}"
        return text


def _read_with_fallback(path: Path) -> str:
    """尝试多种编码读取文件。"""
    encodings = ["utf-8-sig", "utf-8", "gbk", "latin-1"]
    for enc in encodings:
        try:
            return path.read_text(encoding=enc)
        except (UnicodeDecodeError, UnicodeError):
            continue
    raise ValueError(f"无法识别文件编码: {path}")


# ============================================================
# PDF 提取器
# ============================================================

class PdfExtractor(BaseExtractor):
    """
    PDF 提取器。
    逐页提取文本，保留页码顺序，添加 ## Page N 标记。
    如果某页无可提取文本，插入占位提示。
    """

    def extract(self, path: Path) -> str:
        try:
            from pypdf import PdfReader
        except ImportError:
            raise ImportError(
                "需要安装 pypdf 库来提取 PDF。请运行: pip install pypdf"
            )

        reader = PdfReader(str(path))
        title = path.stem
        lines: list[str] = [f"# {title}", ""]

        for i, page in enumerate(reader.pages, start=1):
            text = page.extract_text()
            lines.append(f"## Page {i}")
            lines.append("")

            if text and text.strip():
                # 保留段落结构
                paragraphs = _clean_text(text)
                lines.append(paragraphs)
            else:
                lines.append(
                    "[This page may contain images, charts, "
                    "or non-extractable content.]"
                )
            lines.append("")

        return "\n".join(lines)


# ============================================================
# PPTX 提取器
# ============================================================

class PptxExtractor(BaseExtractor):
    """
    PPTX 提取器。
    逐幻灯片提取：标题 → 文本框/占位符正文 → 演讲者备注。
    每页添加 ## Slide N 标记。
    如果某页无文本，插入占位提示。
    """

    def extract(self, path: Path) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError(
                "需要安装 python-pptx 库来提取 PPTX。请运行: pip install python-pptx"
            )

        prs = Presentation(str(path))
        title = path.stem
        lines: list[str] = [f"# {title}", ""]

        for i, slide in enumerate(prs.slides, start=1):
            lines.append(f"## Slide {i}")
            lines.append("")

            # 提取标题
            slide_title = _extract_slide_title(slide)
            if slide_title:
                lines.append(f"Title: {slide_title}")
                lines.append("")

            # 提取文本框和占位符正文
            body_texts = _extract_slide_body(slide)
            if body_texts:
                lines.append(body_texts)
                lines.append("")

            # 提取演讲者备注
            notes_text = _extract_slide_notes(slide)
            if notes_text:
                lines.append(f"### Speaker Notes")
                lines.append("")
                lines.append(notes_text)
                lines.append("")

            # 如果这一页没有任何文本内容
            if not slide_title and not body_texts and not notes_text:
                lines.append(
                    "[This slide may contain images, diagrams, "
                    "or non-extractable content.]"
                )
                lines.append("")

        return "\n".join(lines)


def _extract_slide_title(slide) -> str | None:
    """提取幻灯片标题。"""
    from pptx import Presentation

    if slide.shapes.title and slide.shapes.title.text:
        return slide.shapes.title.text.strip()
    # 回退：找第一个有文本的形状作为标题
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text.strip():
            # 优先取顶部的大号文本（模拟标题）
            return shape.text.strip().split("\n")[0]
    return None


def _extract_slide_body(slide) -> str:
    """提取幻灯片正文（排除已作为标题的内容）。"""
    from pptx import Presentation

    title_shape = slide.shapes.title
    title_text = title_shape.text.strip() if title_shape and title_shape.text else ""

    parts: list[str] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text.strip()
        if not text:
            continue
        # 跳过已作为标题的内容
        if text == title_text:
            continue
        # 跳过与标题完全一致的第一行
        lines_in_shape = text.split("\n")
        if lines_in_shape[0].strip() == title_text:
            text = "\n".join(lines_in_shape[1:]).strip()
            if not text:
                continue
        parts.append(text)

    return "\n\n".join(parts)


def _extract_slide_notes(slide) -> str | None:
    """提取幻灯片演讲者备注。"""
    try:
        notes = slide.notes_slide
        if notes and notes.notes_text_frame:
            text = notes.notes_text_frame.text.strip()
            if text:
                return text
    except Exception:
        pass
    return None


def _clean_text(text: str) -> str:
    """清理提取文本：压缩多余空白，保留段落。"""
    # 统一换行
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    # 按段落处理
    paragraphs = []
    for para in text.split("\n"):
        stripped = para.strip()
        if stripped:
            paragraphs.append(stripped)
    # 用双换行连接段落
    return "\n\n".join(paragraphs)


# ============================================================
# 主入口
# ============================================================

def main():
    parser = argparse.ArgumentParser(
        description="从 PDF/PPTX/MD/TXT 提取文本，生成统一 Markdown。"
    )
    parser.add_argument(
        "--input", "-i", required=True,
        help="输入文件路径（.pdf / .pptx / .md / .txt）",
    )
    parser.add_argument(
        "--output", "-o", default=None,
        help="输出 .md 路径（默认: extracted/<stem>.md）",
    )
    args = parser.parse_args()

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"[extract] 错误：输入文件不存在: {input_path}")
        sys.exit(1)

    # 解析输出路径
    if args.output:
        output_path = Path(args.output)
    else:
        EXTRACTED_DIR.mkdir(parents=True, exist_ok=True)
        output_path = EXTRACTED_DIR / f"{input_path.stem}.md"

    print(f"[extract] 输入: {input_path}")
    print(f"[extract] 格式: {input_path.suffix.upper()}")

    # 创建提取器并执行
    try:
        extractor = BaseExtractor.for_path(input_path)
        markdown = extractor.extract(input_path)
    except Exception as e:
        print(f"[extract] 提取失败: {e}")
        sys.exit(1)

    # 写出结果
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(markdown, encoding="utf-8")
    print(f"[extract] 输出: {output_path}")
    print(f"[extract] 行数: {len(markdown.splitlines())}")
    print(f"[extract] 完成")


if __name__ == "__main__":
    main()
